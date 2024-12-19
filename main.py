import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches
import io
import tempfile

st.title("📄 PDF를 PPTX로 변환하기 🖼️")

uploaded_file = st.file_uploader("PDF 파일을 선택하세요", type="pdf")

if uploaded_file is not None:
    st.success("PDF 파일이 성공적으로 업로드되었습니다!")

    # PDF를 이미지로 변환
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_file_path = tmp_file.name

    images = convert_from_bytes(uploaded_file.getvalue())

    # PPTX 생성
    prs = Presentation()
    for image in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
        
        # PDF 페이지 크기에 맞춰 슬라이드 크기 조정
        prs.slide_width = Inches(image.width / 96)
        prs.slide_height = Inches(image.height / 96)
        
        # 이미지를 슬라이드에 추가
        left = top = Inches(0)
        pic = slide.shapes.add_picture(io.BytesIO(image.tobytes()), left, top, width=prs.slide_width, height=prs.slide_height)

    # PPTX 파일 저장
    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)

    # 다운로드 버튼 생성
    st.download_button(
        label="PPTX 다운로드",
        data=pptx_file,
        file_name="converted_presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

st.info("PDF 파일을 업로드하면 PPTX로 변환됩니다. 변환된 파일은 다운로드 버튼을 통해 받을 수 있습니다.")
